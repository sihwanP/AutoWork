# -------------------------------------- 25.11.17 (월)---------------------------------------------------------
# 폴더와 파일 다루기
# from openpyxl import load_workbook
# from openpyxl import Workbook
# from openpyxl.styles import Font, PatternFill, Side, Border, Alignment
# from datetime import datetime
# from openpyxl import load_workbook

'''
# if not os.path.exists("C:/Sihwan/Book"):
#   print("폴더 없음")
# else:
#   print("폴더 있음")

# lists = os.listdir("C:/Sihwan/code/excel")
# print(lists)

# 폴더이름 변경
# if os.path.exists("C:/Sihwan/code/excel"):
#   os.rename("C:/Sihwan/code/excel", "C:/Sihwan/code/Newexcel") #원본폴더이름, 변경될 폴더 이름

#폴더 복사
# path_from = "C:/Sihwan/code/Newexcel"
# path_to = "C:/Newexcel"
# if not os.path.exists(path_to):
#   shutil.copytree(path_from, path_to)


#파일 내용 읽기.
# 파일 내용을 기록,수정,추가를 하면 .close()로 닫아야 한다.
file = open("example.txt", "r", encoding="utf-8")
# r:읽기모드  w:기록하기(저장) | a: 내용추가(수정) | x: 해당하는 파일이 없으면 만든다.(덮어씌우기)
content = file.read()
file.close()
print(content)

# with는 자동으로 .close()가 된다. (윗 코드랑 같은 의미)
with open("example.txt", "w", encoding="utf-8") as file:
  file.write("홍길동\n안녕하세요")


with open("example.txt", "r", encoding="utf-8") as file:
  line1 = file.readline()
  line2 = file.readline()

print(line1, line2)

# 공백 제거 -> 전부 출력
with open("example.txt", "r", encoding="utf-8") as file:
  line = file.readline()

  while line:
    print(line.strip()) # .strip() : 공백을 전부 제거
    line = file.readline()

# readlines의 s 같은건 가급적 사용 금지!!
with open("example.txt", "r", encoding="utf-8") as file:
  line = file.readlines()
  print(line)

with open("example.txt", "r", encoding="utf-8") as file:
  line = file.readline()
  print(line)

with open("example.csv", "w", encoding="cp949", newline="") as file:
  csv_writer = csv.writer(file)
  csv_writer.writerow(["이름","나이","직업"])
  csv_writer.writerow(["홍길동","29","취준생"])
  csv_writer.writerow(["박시환","30","직장인"])
  csv_writer.writerow(["희야","34","직장인"])
  csv_writer.writerow(["날좀","25","프리"])
  csv_writer.writerow(["바라봐","30","직장인"])
'''

'''엑셀 연동
from openpyxl import Workbook

wb = Workbook()
ws = wb.active
ws.title = "수강생 정보"

# ws["A1"] = "이철수"
# wb.save("수강생 리스트.xlsx")
# wb.close()

column = ["번호", "이름", "과목"]
ws.append(column)
row = [[1,"이철수","수학"],[2, "빛나리", "영어"],[1,"홍길동","수학"]]
for data in row:
  ws.append(data)
# row = [1, "이철수", "수학"]
# ws.append(row)

# 시트 자동 생성
# wb.create_sheet("중간 평가")
# wb.create_sheet("기말 평가")
wb.save("수강생_리스트.xlsx")
wb.close()

'''

'''
wb = load_workbook(filename="월별구매고객리스트.xlsx")
ws = wb["10월"]

new_rows = list(ws.rows)[2:]

for row in new_rows:
  row_values = [cell.value for cell in row]
  print(row_values)
'''

# wb = Workbook()
# ws = wb.active

'''
# 엑셀 백지 상태( 새문서 상태)
cell = ws['A1']
cell.value = "Hello World"

cell.font = Font(color = 'FF0000', italic = True, bold = True, size = 20)

ws.column_dimensions['A'].width = 50
ws.row_dimensions[1].height = 50

yello_fill = PatternFill(start_color='FFFF00', end_color='FFFF00', fill_type='solid')
cell.fill = yello_fill
# 개별 선의 스타일을 먼저 정의한다 > 적용할 부분 설정
thin_side = Side(style = 'thin')
cell.border = Border(left = thin_side, right = thin_side, top = thin_side, bottom = thin_side,)

cell.alignment = Alignment(horizontal = 'center', vertical = 'center')

wb.save('엑셀 서식.xlsx')
'''

'''
ws['A1'] = 1234567.890123
ws['A1'].number_format = "#,##0.00"

ws['D1'] = 1234567.890123
ws['D1'].number_format = "#,##0"

ws['B1'] = 1234567.890123
ws['B1'].number_format = "0.00%"

ws['A1'] = '2025-11-18'
ws['A1'].number_format = 'yyyy-MM-DD'

wb.save('엑셀 서식.xlsx')
'''

# -------------------------------------- 25.11.18 (화)---------------------------------------------------------
'''
# 월별구매고객리스트 중에 10월 시트
wb = load_workbook(filename='월별구매고객리스트.xlsx', data_only = True)
ws = wb['10월']
new_rows = list(ws.rows)[2:]

# new_rows 리스트에 있는 행들을 하나씩 가져와서 반복 작업 진행
for row in new_rows:
  # 조건: 현재 행에서 왼쪽에서 6번째 셀(칸)을 의미한다
  if row[5].value == '신규':
    #            각 셀에서 값만 추출
    row_value = [cell.value for cell in row]
                            # 현재 행의 모든 셀(칸)을 처음부터 마지막까지 순회
    print(row_value)

wb.save(filename='제품별신규고객리스트.xlsx')
'''
# ----------------------------------(제품별신규고객리스트_추출본)-----------------------------------------------
'''
from openpyxl import load_workbook, Workbook  # 👈 Workbook 추가 불러오기

# 1. 기존 파일 열기
wb_source = load_workbook(filename='월별구매고객리스트.xlsx', data_only=True)
ws_source = wb_source['10월']
new_rows = list(ws_source.rows)[2:]

# 2. 새 워크북 및 시트 생성
wb_new = Workbook()
# 새로 만들어진 기본 시트 (Sheet)를 선택합니다.
ws_new = wb_new.active
ws_new.title = '신규 고객 리스트'  # 시트 이름 변경

# 3. 헤더(제목 줄) 복사 및 추가
# 원본 시트의 첫 번째 행 (헤더)을 가져와 새 시트에 추가합니다.
header = [cell.value for cell in list(ws_source.rows)[0]]
ws_new.append(header)

# 4. 조건 확인 및 새 시트에 데이터 추가
for row in new_rows:
    # 조건: 6번째 셀의 값이 '신규'인 경우
    if row[5].value == '신규':
        # 현재 행의 모든 셀에서 값만 추출
        row_values = [cell.value for cell in row]

        # 👈 추출된 값들(리스트)을 새 시트의 새로운 행으로 추가
        ws_new.append(row_values)

    # 5. 새 워크북 저장
# 새로운 워크북 객체 (wb_new)를 저장합니다.
wb_new.save(filename='제품별신규고객리스트_추출본.xlsx')

print("신규 고객 정보만 담긴 '제품별신규고객리스트_추출본.xlsx' 파일이 생성되었습니다.")
'''

# --------------------------------(제품별 신규 고객 리스트.xlsx)------------------------------------------------
'''
wb_source = load_workbook(filename='월별구매고객리스트.xlsx', data_only = True)
month = ['10월', '11월', '12월']
for month  in month:
  ws = wb_source[month] # 해당 월의 데이터 가져오기
  new_rows = list(ws.rows)[2:]
  for row in new_rows:
    row_value = [cell.value for cell in row]
    print(row_value)

wb_source.save(filename='제품별 신규 고객 리스트.xlsx')
'''

# --------------------------------------------(fax)-----------------------------------------------------------
'''
from openpyxl import Workbook, load_workbook
import os
import re  # Regular Expression : 정규 표현식

# 새로운 워크북과 워크시트를 만든다.
new_wb = Workbook()
new_ws = new_wb.active
new_ws.title = '수신내역' # 워크시트의 제목 설정

# 첫 번째 행에 각각의 열을 의미하는 대표 제목
new_ws.append(['수신시간', '발신번호', '페이지수', '용량'])

for filename in os.listdir('./fax'):
  # fax 폴더에 있는 모든 파일의 목록을 가져온다

  wb = load_workbook('./fax/' + filename)
  # 현재 파일을 열어서 wb 객체를 만든다

  ws = wb.active
  # 열린 워크북의 활성화된 워크시트를 의미한다.

  for row in ws.iter_rows(min_row = 2, values_only = True):
    # 첫 번째 행 제목(제목 행)을 제외 (두번째 행 부터)
    new_ws.append(row)
    # 새로운 워크시트에 (새로운 워크북)에 추가한다

new_wb.save('10월팩스_수신내역.xlsx')
'''

# --------------------------------------------(docx)-----------------------------------------------------------
'''
# from docx import Document #문서
# from docx.enum.text import WD_ALIGN_PARAGRAPH
# import docx2txt
# import os

# doc = Document()
#
# title = doc.add_heading('제목을 이곳에 작성합니다', level=1) #level=0 =>숫자가 커질수록 제목 종류가 달라짐.(0~9까지)
# title.alignment = WD_ALIGN_PARAGRAPH.CENTER
#
# p = doc.add_paragraph('여기는 단락입니다.')

# table = doc.add_table(rows = 3, cols = 3)
# table.style = 'Table Grid'
#
# hdr_cells = table.rows[0].cells
# hdr_cells[0].text = '헤더 1'
# hdr_cells[1].text = '헤더 2'
# hdr_cells[2].text = '헤더 3'
# for i in range(1, 3):
#   row_cells = table.rows[i].cells
#   row_cells[0].text = f'행 {i}, 열 1'
#   row_cells[1].text = f'행 {i}, 열 2'
#   row_cells[2].text = f'행 {i}, 열 3'
# doc.save('example.docx')
#
#
# save_path = './word_savepath/'
#
# if not os.path.exists(save_path):
#   os.mkdir(save_path)
# docx2txt.process('example_image.docx', save_path)
#
# # 비파괴적인(non-destructive)
'''
# --------------------------------------------(PPT)-----------------------------------------------------------
'''
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "제목!!"

subtitle = slide.placeholders[1]
subtitle.text = '부제목'

left = Inches(1)
top = Inches(2)
width = Inches(5)
height = Inches(1.5)

textbox = slide.shapes.add_textbox(left = left, top = top, width = width, height = height)
frame = textbox.text_frame
p = frame.add_paragraph()
p.text = '새로운 텍스트 박스에 추가된 문자열'

# for i, placeholder in enumerate(slide.placeholders):
#   print(f'Placeholder {i}: {placeholder.name}')

# for i, layout in enumerate(prs.slide_layouts):
#   print(f'Layout {i}: {layout.name}')

# slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(slide_layout)

prs.save('Presentation.pptx')
'''


# PPT 하위 만들기
from pptx import Presentation

prs = Presentation()
slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = '파이썬의 장점'

# ➊ Placeholder 1의 TextFrame 객체에 리스트의 첫 번째 항목 입력
tf = slide.placeholders[1].text_frame
tf.text = '쉬운 사용법'

# ➋ 첫 번째 항목에 대한 하위 항목 추가
p = tf.add_paragraph()
p.text = '직관적인 문법'
p.level = 1

# ➌ 리스트의 두 번째 항목 추가
p = tf.add_paragraph()
p.text = '높은 생산성'
p.level = 0

# ➍ 두 번째 항목에 대한 하위 항목 추가
p = tf.add_paragraph()
p.text = '빠른 개발 속도'
p.level = 1

# ➎ 리스트의 세 번째 항목 추가
p = tf.add_paragraph()
p.text = '다양한 라이브러리와 프레임워크'
p.level = 0

# ➏ 세 번째 항목에 대한 하위 항목 추가
p = tf.add_paragraph()
p.text = '머신 러닝, 웹 개발 등에 유용'
p.level = 1

prs.save('presentation_with_list.pptx')










